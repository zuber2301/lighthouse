from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import os

presentation_path = "presentation/Product_Presentation.pptx"

# Simple slide definitions
slides = [
    {
        "title": "Product Name — Unified Recognition & Engagement Platform",
        "bullets": [
            "Customer-facing product overview and user workflows",
            "Presenter: [Your Name] — Date: [Presentation Date]"
        ],
        "notes": "Title slide. Introduce yourself and set expectations for the 20-30 minute demo/presentation."
    },
    {"title": "Agenda", "bullets": ["Product overview and value proposition", "Features and differentiators", "Workflows: how customers use the product", "Personas and benefits", "Getting Started and support", "Commercial terms and next steps"], "notes": "Run through the agenda briefly to set expectations."},
    {"title": "Product Overview", "bullets": ["Centralized SaaS platform for employee recognition, rewards, and multi-tenant management", "Purpose: increase engagement, simplify reward operations, provide measurable insights", "Key outcomes: adoption, analytics, reduced admin effort"], "notes": "Explain the core problem the product solves."},
    {"title": "Value Proposition", "bullets": ["HR & People Ops: centralized controls and audit-ready reporting", "Managers: fast recognition flows and team engagement visibility", "Finance & Execs: business metrics (MRR impact, adoption) for decision-making"], "notes": "Tailor examples to the customer's org."},
    {"title": "Key Capabilities", "bullets": ["Recognition engine: badges, approvals, configurable flows", "Rewards & redemption: catalogs, budget controls, payout reporting", "Analytics: aggregated MRR, active users, top tenants, trends", "Tenant management: RBAC and lifecycle tools", "Integrations & security: SSO, Slack, HRIS, encryption"], "notes": "Highlight capabilities most relevant to the customer."},
    {"title": "Recognition — Feature Deep Dive", "bullets": ["Create & launch: templates, categories, approval settings", "Channels: web UI, Slack, email notifications", "Governance: budget caps, approval routing, audit logs"], "notes": "Show real-world examples during demo."},
    {"title": "Rewards — Feature Deep Dive", "bullets": ["Catalogs: digital & physical rewards, points model", "Redemption: seamless user experience, automated fulfillment hooks", "Finance: exportable payout reports & accounting integrations"], "notes": "Emphasize finance controls if relevant."},
    {"title": "Analytics — Feature Deep Dive", "bullets": ["Dashboards: Aggregated MRR, growth %, active users, uptime, top tenants", "Exploration: tenant/team/user drilldowns, cohorts, time-series", "Actions: alerts, scheduled reports, CSV/BI exports"], "notes": "Offer examples of KPIs to track during the pilot."},
    {"title": "Workflow — Day-to-Day (Recognition)", "bullets": ["Manager or peer submits recognition", "Optional approval triggers", "Points assigned and available for redemption", "User redeems; finance receives reports"], "notes": "Walk through a short demo flow here."},
    {"title": "Workflow — Admin (Tenant Management)", "bullets": ["Create tenant, configure branding and roles", "Connect SSO and HRIS, import users", "Configure recognition catalog and budgets", "Monitor adoption and adjust settings"], "notes": "Describe minimal admin effort required."},
    {"title": "Getting Started — Customer Usage", "bullets": ["First 48 hours: create tenant, invite admins, set one badge, run a pilot", "First 2 weeks: pilot with one team, configure SSO + HRIS sync", "Ongoing: monthly reports, champion program, campaigns"], "notes": "Set realistic expectations for a pilot rollout."},
    {"title": "Personas", "bullets": ["HR Admin: onboards tenants, sets policies, runs reports", "People Manager: gives recognition, tracks team health", "Employees: receive recognition, redeem rewards", "CTO/Platform Admin: integrations, security, uptime"], "notes": "Relate personas to the stakeholders in the customer's org."},
    {"title": "Support & Success", "bullets": ["Onboarding: live kickoff, admin training, runbook", "Help: in-app help, knowledge base, Slack/email support", "Success metrics: adoption rate, recognitions/month, redemption rate"], "notes": "Clarify channels and SLAs for support."},
    {"title": "Commercial", "bullets": ["Tiers: Pilot, Professional, Enterprise (SLA & integrations)", "Pricing model: per-tenant base + active-user bands; add-ons for custom work", "Trial: 2-week pilot proposal with success criteria"], "notes": "Be ready to discuss pricing range and pilot terms."},
    {"title": "Case Study / Example", "bullets": ["Customer profile: industry & size", "Problem solved: fragmented recognition + low adoption", "Result: measurable uplift in recognitions and engagement KPIs"], "notes": "Use a short, concrete case study if available."},
    {"title": "Next Steps / Call to Action", "bullets": ["Request: scope a 2-week pilot with [Customer Name]", "Deliverables: pilot plan, success criteria, kickoff date", "Contact: [Your Name] — [Email / Phone]"], "notes": "End with a clear ask and next steps."},
    {"title": "Appendix / Demo Notes", "bullets": ["Demo focus: Dashboard (Aggregated MRR, active users), recognition flow, redemption flow", "Data required for demo: sample tenant, sample recognitions, SSO test user", "Speaker tips: highlight ease-of-use and measurable outcomes"], "notes": "Keep appendix slides for Q&A or demo stage directions."}
]

prs = Presentation()

# Prepare assets directory and simple placeholder images (logo, icon, screenshot)
assets_dir = os.path.join(os.path.dirname(presentation_path), "assets")
os.makedirs(assets_dir, exist_ok=True)

# Helper to make a simple placeholder image with text
def make_placeholder(path, size, color, text):
    img = Image.new("RGB", size, color)
    draw = ImageDraw.Draw(img)
    try:
        f = ImageFont.truetype("DejaVuSans-Bold.ttf", 20)
    except Exception:
        f = ImageFont.load_default()
    try:
        bbox = draw.textbbox((0, 0), text, font=f)
        w = bbox[2] - bbox[0]
        h = bbox[3] - bbox[1]
    except Exception:
        w, h = f.getsize(text)
    draw.text(((size[0]-w)/2, (size[1]-h)/2), text, fill=(255,255,255), font=f)
    img.save(path)

logo_path = os.path.join(assets_dir, "logo.png")
icon_path = os.path.join(assets_dir, "icon.png")
ss_path = os.path.join(assets_dir, "screenshot.png")
if not os.path.exists(logo_path):
    make_placeholder(logo_path, (300, 80), (30, 90, 180), "LOGO")
if not os.path.exists(icon_path):
    make_placeholder(icon_path, (120, 120), (80, 160, 80), "ICON")
if not os.path.exists(ss_path):
    make_placeholder(ss_path, (1280, 720), (70, 70, 70), "SAMPLE SCREENSHOT")

# small helper to add logo on top-right of a slide
def add_logo(slide):
    try:
        slide.shapes.add_picture(logo_path, prs.slide_width - Inches(1.6), Inches(0.2), width=Inches(1.4))
    except Exception:
        pass

# Title slide (use layout 0)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = slides[0]["title"]
subtitle = slide.placeholders[1]
subtitle.text = "\n".join(slides[0]["bullets"])
slide.notes_slide.notes_text_frame.text = slides[0]["notes"]
add_logo(slide)

# Other slides
for s in slides[1:]:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = s["title"]
    # increase title font size where possible
    try:
        title.text_frame.paragraphs[0].font.size = Pt(28)
    except Exception:
        pass
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, b in enumerate(s["bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
            p.font.size = Pt(14)
        else:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(12)
            p.level = 0
    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = s.get("notes", "")
    add_logo(slide)

    # Add visuals on select slides
    if s["title"] in ("Product Overview", "Appendix / Demo Notes"):
        try:
            slide.shapes.add_picture(ss_path, Inches(0.5), Inches(3.0), width=prs.slide_width - Inches(1.0))
        except Exception:
            pass
    # Add a small icon next to title for feature deep dives
    if "Deep Dive" in s["title"] or "Key Capabilities" == s["title"]:
        try:
            slide.shapes.add_picture(icon_path, Inches(0.4), Inches(0.6), width=Inches(0.6))
        except Exception:
            pass

# Ensure directory exists
import os
os.makedirs(os.path.dirname(presentation_path), exist_ok=True)
prs.save(presentation_path)
print(f"Saved presentation to: {presentation_path}")

